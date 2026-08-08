"""
Unit tests for all extractor functions.

Task 4.9 — Tests for PDF scanned-page skipping, DOCX/PPTX/plain-text
extraction, and unhandled-exception propagation through detect_and_route.

Validates: Requirements 3.2, 3.6
"""

import io
import sys
import os
import logging
from unittest.mock import patch

# ---------------------------------------------------------------------------
# Path setup — makes the file runnable from any working directory
# ---------------------------------------------------------------------------
_BACKEND_DIR = os.path.join(os.path.dirname(__file__), "..")
if _BACKEND_DIR not in sys.path:
    sys.path.insert(0, os.path.abspath(_BACKEND_DIR))

import pytest
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.util import Inches

from services.extraction.pdf import extract_pdf
from services.extraction.docx import extract_docx
from services.extraction.pptx import extract_pptx
from services.extraction.plain_text import extract_plain_text
from services.extraction.router import detect_and_route


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_pdf(page_texts: list[str]) -> bytes:
    """
    Create an in-memory PDF with one page per entry in *page_texts*.
    Each entry is inserted as a text annotation on the page.
    """
    doc = fitz.open()
    for text in page_texts:
        page = doc.new_page()
        page.insert_text((72, 72), text)
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    return buf.getvalue()


def _make_docx(paragraphs: list[str]) -> bytes:
    """Create an in-memory DOCX with the given paragraph texts."""
    document = DocxDocument()
    for text in paragraphs:
        document.add_paragraph(text)
    buf = io.BytesIO()
    document.save(buf)
    return buf.getvalue()


def _make_pptx(slide_texts: list[str]) -> bytes:
    """Create an in-memory PPTX with one slide per entry in *slide_texts*."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # completely blank layout
    for text in slide_texts:
        slide = prs.slides.add_slide(blank_layout)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        txBox.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# PDF tests
# ---------------------------------------------------------------------------

class TestPdfExtractor:

    def test_scanned_page_skipped_only_normal_page_returned(self, caplog):
        """
        Validates: Requirement 3.2

        A PDF with a near-empty page (< 50 chars) alongside a normal page
        should yield exactly one segment (the normal page) and log a warning
        about the skipped page.
        """
        short_text = "Hi"  # well under 50 chars
        normal_text = "A" * 60  # 60 chars — above the 50-char threshold

        pdf_bytes = _make_pdf([short_text, normal_text])

        with caplog.at_level(logging.WARNING, logger="services.extraction.pdf"):
            segments = extract_pdf(pdf_bytes)

        # Only one segment — the normal page
        assert len(segments) == 1
        assert segments[0].location == 2

        # Warning should mention page 1
        warning_messages = [r.message for r in caplog.records if r.levelno == logging.WARNING]
        assert any("1" in msg for msg in warning_messages), (
            f"Expected a warning mentioning page 1, got: {warning_messages}"
        )

    def test_all_normal_pages_both_returned(self):
        """Both pages with ≥ 50 chars → 2 segments with locations [1, 2]."""
        text_a = "B" * 60
        text_b = "C" * 60

        pdf_bytes = _make_pdf([text_a, text_b])
        segments = extract_pdf(pdf_bytes)

        assert len(segments) == 2
        assert segments[0].location == 1
        assert segments[1].location == 2


# ---------------------------------------------------------------------------
# DOCX tests
# ---------------------------------------------------------------------------

class TestDocxExtractor:

    def test_basic_three_paragraphs(self):
        """3 paragraphs → 3 segments, all with location 1 (ceil(1/10)=…=ceil(3/10)=1)."""
        paragraphs = [
            "First paragraph with some text.",
            "Second paragraph with some text.",
            "Third paragraph with some text.",
        ]
        docx_bytes = _make_docx(paragraphs)
        segments = extract_docx(docx_bytes)

        assert len(segments) == 3
        for seg in segments:
            assert seg.location == 1

    def test_location_formula_eleven_paragraphs(self):
        """
        11 non-empty paragraphs:
        - paragraph 10 → location = ceil(10/10) = 1
        - paragraph 11 → location = ceil(11/10) = 2
        """
        paragraphs = [f"Paragraph {i} has enough content here." for i in range(1, 12)]
        docx_bytes = _make_docx(paragraphs)
        segments = extract_docx(docx_bytes)

        assert len(segments) == 11
        assert segments[9].location == 1   # 10th segment (0-indexed: 9)
        assert segments[10].location == 2  # 11th segment (0-indexed: 10)


# ---------------------------------------------------------------------------
# PPTX tests
# ---------------------------------------------------------------------------

class TestPptxExtractor:

    def test_basic_two_slides(self):
        """2 slides with text → 2 segments with locations [1, 2]."""
        slide_texts = [
            "Slide one content with some text here.",
            "Slide two content with some text here.",
        ]
        pptx_bytes = _make_pptx(slide_texts)
        segments = extract_pptx(pptx_bytes)

        assert len(segments) == 2
        assert segments[0].location == 1
        assert segments[1].location == 2


# ---------------------------------------------------------------------------
# Plain-text / Markdown tests
# ---------------------------------------------------------------------------

class TestPlainTextExtractor:

    def test_plain_text_single_segment_no_location(self):
        """Plain text → exactly 1 segment, location is None."""
        segments = extract_plain_text(b"Hello world, this is a test")

        assert len(segments) == 1
        assert segments[0].location is None

    def test_markdown_single_segment_no_location(self):
        """Markdown content → exactly 1 segment, location is None."""
        segments = extract_plain_text(b"# Heading\n\nSome markdown text")

        assert len(segments) == 1
        assert segments[0].location is None


# ---------------------------------------------------------------------------
# Unhandled extraction exception propagation — Requirement 3.6
# ---------------------------------------------------------------------------

class TestUnhandledExtractionException:

    def test_unhandled_runtime_error_propagates_from_router(self):
        """
        Validates: Requirement 3.6

        When the underlying extractor raises an unhandled RuntimeError,
        detect_and_route must re-raise it (the HTTP 500 conversion is done
        at the upload route level, not inside the router).
        """
        pdf_bytes = b"%PDF-1.4 " + b"x" * 20  # valid magic bytes → routed to pdf

        with patch(
            "services.extraction.router.extract_pdf",
            side_effect=RuntimeError("boom"),
        ):
            with pytest.raises(RuntimeError, match="boom"):
                detect_and_route(pdf_bytes, "doc.pdf", "application/pdf")
