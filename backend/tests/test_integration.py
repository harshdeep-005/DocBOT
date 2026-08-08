"""
tests/test_integration.py — Backend integration tests (mocked at the Gemini
                             network boundary).

All tests exercise the full pipeline through the FastAPI ASGI app:
  POST /upload → ingest document
  POST /ask    → RAG query

Gemini API calls (embed_content, generate_content) are mocked so no real API
key or network access is required.

Requirements validated: 1.7, 3.1, 3.4, 3.5, 7.1
"""

import io
import re

import fitz  # PyMuPDF
import pytest
import pytest_asyncio
from pptx import Presentation
from pptx.util import Inches

# conftest.py provides: test_client, mock_gemini, test_env


# ── Helpers — build minimal synthetic files in memory ─────────────────────────

def _make_pdf_bytes(page_texts: list[str] | None = None) -> bytes:
    """
    Build a minimal valid PDF with one page per entry in page_texts.
    Defaults to a single page with sample content if page_texts is None.
    """
    if page_texts is None:
        page_texts = ["This is a sample page about artificial intelligence and machine learning."]

    doc = fitz.open()
    for text in page_texts:
        page = doc.new_page()
        page.insert_text((72, 72), text)
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    return buf.getvalue()


def _make_pptx_bytes(slide_texts: list[str] | None = None) -> bytes:
    """
    Build a minimal valid PPTX with one slide per entry in slide_texts.
    Defaults to a single slide with sample content if slide_texts is None.
    """
    if slide_texts is None:
        slide_texts = ["This slide covers neural networks and deep learning."]

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank layout

    for text in slide_texts:
        slide = prs.slides.add_slide(blank_layout)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
        tf = txBox.text_frame
        tf.text = text

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_txt_bytes(content: str | None = None) -> bytes:
    """Return UTF-8-encoded plain text bytes."""
    if content is None:
        content = "This is a plain text document. It contains information about Python programming."
    return content.encode("utf-8")


# ── Tests ──────────────────────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_upload_pdf_and_ask_question(test_client, mock_gemini, test_env):
    """
    Validates: Requirements 3.1, 3.4, 7.1

    Upload a real minimal PDF:
      - Response must include a non-empty document_id and chunk_count > 0.
    Then ask a question:
      - Answer must be non-empty.
      - Citations must use the "Page N" format (at least one citation).
    """
    pdf_bytes = _make_pdf_bytes(
        ["Page one content about machine learning and AI systems."]
    )

    # ── Upload ────────────────────────────────────────────────────────────────
    upload_response = await test_client.post(
        "/upload/",
        files={"file": ("test.pdf", pdf_bytes, "application/pdf")},
    )
    assert upload_response.status_code == 200, (
        f"Upload failed: {upload_response.status_code} {upload_response.text}"
    )
    upload_data = upload_response.json()
    assert upload_data["document_id"], "document_id must be a non-empty string"
    assert upload_data["chunk_count"] > 0, (
        f"chunk_count must be > 0, got {upload_data['chunk_count']}"
    )

    document_id = upload_data["document_id"]

    # ── Ask ───────────────────────────────────────────────────────────────────
    ask_response = await test_client.post(
        "/ask/",
        json={"question": "What is discussed in this document?", "document_id": document_id},
    )
    assert ask_response.status_code == 200, (
        f"Ask failed: {ask_response.status_code} {ask_response.text}"
    )
    ask_data = ask_response.json()

    assert ask_data["answer"], "Answer must be non-empty"

    # At least one citation should use "Page N" format (Requirement 7.1)
    citations = ask_data["citations"]
    assert len(citations) >= 1, "Expected at least one citation for a PDF upload"
    for citation in citations:
        assert re.match(r"^Page \d+$", citation["label"]), (
            f"PDF citation must be 'Page N' format, got: {citation['label']!r}"
        )


@pytest.mark.asyncio
async def test_upload_pptx_citations(test_client, mock_gemini, test_env):
    """
    Validates: Requirements 3.5, 7.1

    Upload a real minimal PPTX, then ask a question.
    All citations must use the "Slide N" format.
    """
    pptx_bytes = _make_pptx_bytes(
        ["Slide one about neural networks.", "Slide two about deep learning."]
    )

    # ── Upload ────────────────────────────────────────────────────────────────
    upload_response = await test_client.post(
        "/upload/",
        files={
            "file": (
                "test.pptx",
                pptx_bytes,
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )
    assert upload_response.status_code == 200, (
        f"PPTX upload failed: {upload_response.status_code} {upload_response.text}"
    )
    upload_data = upload_response.json()
    assert upload_data["chunk_count"] > 0

    document_id = upload_data["document_id"]

    # ── Ask ───────────────────────────────────────────────────────────────────
    ask_response = await test_client.post(
        "/ask/",
        json={"question": "What topics are covered?", "document_id": document_id},
    )
    assert ask_response.status_code == 200, (
        f"Ask failed: {ask_response.status_code} {ask_response.text}"
    )
    ask_data = ask_response.json()

    assert ask_data["answer"], "Answer must be non-empty"

    citations = ask_data["citations"]
    assert len(citations) >= 1, "Expected at least one citation for a PPTX upload"
    for citation in citations:
        assert re.match(r"^Slide \d+$", citation["label"]), (
            f"PPTX citation must be 'Slide N' format, got: {citation['label']!r}"
        )


@pytest.mark.asyncio
async def test_upload_txt_no_citations(test_client, mock_gemini, test_env):
    """
    Validates: Requirements 3.1, 7.1

    Upload a plain-text file. When asking a question the response citations
    list must be empty because TXT chunks carry no location metadata.
    """
    txt_bytes = _make_txt_bytes(
        "Plain text content about Python programming and software development."
    )

    # ── Upload ────────────────────────────────────────────────────────────────
    upload_response = await test_client.post(
        "/upload/",
        files={"file": ("test.txt", txt_bytes, "text/plain")},
    )
    assert upload_response.status_code == 200, (
        f"TXT upload failed: {upload_response.status_code} {upload_response.text}"
    )
    upload_data = upload_response.json()
    assert upload_data["chunk_count"] > 0

    document_id = upload_data["document_id"]

    # ── Ask ───────────────────────────────────────────────────────────────────
    ask_response = await test_client.post(
        "/ask/",
        json={"question": "What is this document about?", "document_id": document_id},
    )
    assert ask_response.status_code == 200, (
        f"Ask failed: {ask_response.status_code} {ask_response.text}"
    )
    ask_data = ask_response.json()

    assert ask_data["answer"], "Answer must be non-empty"

    # TXT files must produce no citations (Requirement 7.1)
    assert ask_data["citations"] == [], (
        f"TXT upload must produce an empty citations list, got: {ask_data['citations']}"
    )


@pytest.mark.asyncio
async def test_reupload_clears_state(test_client, mock_gemini, test_env):
    """
    Validates: Requirement 1.7

    Upload a first document, record its document_id, then upload a second
    document. The second document_id must differ from the first, proving that
    each upload produces a fresh document_id (the frontend uses this change
    to clear prior conversation state).
    """
    pdf_bytes_1 = _make_pdf_bytes(["First document content about machine learning."])
    pdf_bytes_2 = _make_pdf_bytes(["Second document content about data science."])

    # ── First upload ──────────────────────────────────────────────────────────
    upload_response_1 = await test_client.post(
        "/upload/",
        files={"file": ("first.pdf", pdf_bytes_1, "application/pdf")},
    )
    assert upload_response_1.status_code == 200, (
        f"First upload failed: {upload_response_1.status_code} {upload_response_1.text}"
    )
    document_id_1 = upload_response_1.json()["document_id"]
    assert document_id_1, "First document_id must be non-empty"

    # ── Second upload ─────────────────────────────────────────────────────────
    upload_response_2 = await test_client.post(
        "/upload/",
        files={"file": ("second.pdf", pdf_bytes_2, "application/pdf")},
    )
    assert upload_response_2.status_code == 200, (
        f"Second upload failed: {upload_response_2.status_code} {upload_response_2.text}"
    )
    document_id_2 = upload_response_2.json()["document_id"]
    assert document_id_2, "Second document_id must be non-empty"

    # Each upload must produce a distinct document_id (Requirement 1.7)
    assert document_id_1 != document_id_2, (
        f"Each upload must produce a different document_id. "
        f"Got the same id for both uploads: {document_id_1!r}"
    )
