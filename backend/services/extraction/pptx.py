"""
services/extraction/pptx.py — PPTX Extractor.

Extracts text from a PPTX file using python-pptx, producing one
TextSegment per slide with a 1-based slide number as the location.
"""

import io
from typing import List

from pptx import Presentation

from models.schemas import TextSegment


def extract_pptx(file_bytes: bytes) -> List[TextSegment]:
    """
    Extract text from a PPTX file, one TextSegment per slide.

    Each segment's `location` is set to the 1-based slide number.
    Text from all shapes on a slide is joined with newlines.

    Args:
        file_bytes: Raw bytes of the PPTX file.

    Returns:
        List of TextSegment objects, one per slide.
    """
    presentation = Presentation(io.BytesIO(file_bytes))
    segments: List[TextSegment] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs)
                    if line:
                        texts.append(line)

        segments.append(TextSegment(
            text="\n".join(texts),
            location=slide_number,
        ))

    return segments
